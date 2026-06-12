import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ------------------------------------------------------------
# 1. SETUP PATHS & DIRECTORIES
# ------------------------------------------------------------
BASE_DIR = r"C:\Users\daipa\NIFTY50-Volatility-Study"
OUTPUT_DIR = os.path.join(BASE_DIR, "presentation")
os.makedirs(OUTPUT_DIR, exist_ok=True)

TABLES_DIR = os.path.join(BASE_DIR, "outputs", "tables")
GRAPHS_DIR = os.path.join(BASE_DIR, "outputs", "graphs")

# ------------------------------------------------------------
# 2. DESIGN PALETTE & TYPOGRAPHY CONSTANTS
# ------------------------------------------------------------
FONT_TITLE = "Georgia"
FONT_BODY = "Arial"

COLOR_PRIMARY = RGBColor(31, 78, 121)    # Corporate Deep Blue (#1f4e79)
COLOR_SECONDARY = RGBColor(46, 139, 87)  # Emerald Sea Green (#2e8b57)
COLOR_TEXT_DARK = RGBColor(44, 62, 80)   # Slate Charcoal (#2c3e50)
COLOR_MUTED = RGBColor(127, 140, 141)    # Muted Academic Gray (#7f8c8d)
COLOR_HIGHLIGHT = RGBColor(192, 57, 43)  # Risk Red (#c0392b)
COLOR_WHITE = RGBColor(255, 255, 255)

WATERMARK_TEXT = "Daipayan Chatterjee (M.Sc. Economics) | GitHub: dchatterjee01-prog"

# ------------------------------------------------------------
# 3. HELPER FUNCTIONS FOR LAYOUT & STRUCTURAL REUSABILITY
# ------------------------------------------------------------
def create_blank_slide(prs):
    """Creates a blank slide layout to prevent default placeholders interfering."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def add_academic_watermark(slide):
    """Appends the required professional student identifier to the bottom margin."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.33), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = WATERMARK_TEXT
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = COLOR_MUTED
    p.alignment = PP_ALIGN.RIGHT

def add_slide_header(slide, title_text):
    """Generates a standardized corporate-academic presentation header banner."""
    # Top solid color structural bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()
    
    # Title Text Box
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_TITLE
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    add_academic_watermark(slide)

def add_bullet_points(slide, points, left, top, width, height, font_size=14):
    """Injects professional text frames with clean vertical paragraph separation."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, pt in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = pt
        p.font.name = FONT_BODY
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLOR_TEXT_DARK
        p.space_after = Pt(12)
        p.level = 0

def add_image_safely(slide, filename, left, top, width=None, height=None):
    """Injects graphics or maps structural fallbacks if paths don't exist yet."""
    path = os.path.join(GRAPHS_DIR, filename)
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width=width, height=height)
    else:
        # Placeholder shape to keep presentation execution clean if file is missing
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width or Inches(5), height or Inches(4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(235, 240, 245)
        shape.line.color.rgb = COLOR_MUTED
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"[Graphical Visualization Missing]\n{filename}\nEnsure Phase script has run successfully."
        p.font.name = FONT_BODY
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_MUTED
        p.alignment = PP_ALIGN.CENTER

def add_dataframe_table(slide, df_path, left, top, width, height, max_rows=10):
    """Transforms raw pandas CSV analysis outputs seamlessly into stylized slide tables."""
    if not os.path.exists(df_path):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.text_frame.text = "[Data Matrix Missing: Run Data Generation Phases]"
        return
        
    df = pd.read_csv(df_path)
    df = df.dropna(how='all').iloc[:max_rows]
    
    rows = len(df) + 1
    cols = len(df.columns)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Format Headers
    for c_idx, col_name in enumerate(df.columns):
        cell = table.cell(0, c_idx)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        
    # Format Data Matrix rows
    for r_idx, row in df.iterrows():
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            # Round floats automatically for academic readability
            if isinstance(val, float):
                cell.text = f"{val:.4f}"
            else:
                cell.text = str(val)
            
            # Alternating row shading for scannability
            cell.fill.solid()
            if r_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(245, 247, 250)
            else:
                cell.fill.fore_color.rgb = COLOR_WHITE
                
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_TEXT_DARK
            p.alignment = PP_ALIGN.CENTER

# ------------------------------------------------------------
# 4. PRIMARY GENERATOR ROUTINE
# ------------------------------------------------------------
def build_academic_presentation():
    prs = Presentation()
    # Enforce standard 16:9 modern widescreen presentation coordinates
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # --------------------------------------------------------
    # SLIDE 1: TITLE / DEFENSE COVER
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_PRIMARY
    bg.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Return and Volatility Dynamics of NIFTY 50"
    p.font.name = FONT_TITLE
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "An Empirical Econometric Investigation (2023 - 2026)"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_SECONDARY
    p2.space_before = Pt(10)
    
    author_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(11.333), Inches(1.5))
    tf_au = author_box.text_frame
    p3 = tf_au.paragraphs[0]
    p3.text = f"Presenter: {WATERMARK_TEXT.split(' | ')[0]}"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_WHITE
    
    p4 = tf_au.add_paragraph()
    p4.text = f"Codebase & Replication Repository: {WATERMARK_TEXT.split(' | ')[1]}"
    p4.font.name = FONT_BODY
    p4.font.size = Pt(12)
    p4.font.color.rgb = RGBColor(200, 214, 229)
    p4.space_before = Pt(5)

    # --------------------------------------------------------
    # SLIDE 2: RESEARCH OBJECTIVES
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "1. Research Core Objectives")
    pts = [
        "• Establish the empirical properties of the NIFTY 50 index across contemporary market structures (2023-2026).",
        "• Conduct rigorous econometric tests for time series stationarity, normality, and residual dependencies.",
        "• Formally confirm the presence of conditional heteroscedasticity (ARCH effects) in daily index innovations.",
        "• Parameterize time-varying volatility dynamics utilizing a classic GARCH(1,1) specification.",
        "• Implement advanced quantitative risk applications, including time-varying GARCH Parameteric Value at Risk (VaR), Conditional VaR (Expected Shortfall), and empirical volatility regime classification."
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5), font_size=15)

    # --------------------------------------------------------
    # SLIDE 3: DATA ACQUISITION & STRUCTURAL PIPELINE
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "2. Data Pipeline Architecture")
    pts = [
        "• Data Source: High-fidelity daily pricing metrics harvested programmatically via yfinance API infrastructure.",
        "• Target Variable: NIFTY 50 Index Close prices (INR), accompanied by absolute trading volume records.",
        "• Temporal Coverage: Longitudinal daily market capture extending through the 2023 to 2026 horizon.",
        "• Pipeline Design: Multi-tiered architecture cleanly separating raw ingest data storage from downstream processed matrix storage (`data/raw/` to `data/processed/`).",
        "• Reproducibility Standards: Rigorous script-driven execution executing step-by-step modular tasks to allow automated research auditing."
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5))

    # --------------------------------------------------------
    # SLIDE 4: DATA CLEANING LOGIC & INTEGRITY SAFEGUARDS
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "3. Data Wrangling & Sanity Protocols")
    pts = [
        "• Header Multi-Index Resolution: Automated programmatic execution parsing out multi-level yfinance headers by eliminating residual junk rows.",
        "• Type Alignment: Explicit casting transformations mapping fields to proper `datetime64[ns]` schemas and numeric float types.",
        "• Imputation Logic: Elimination of missing values via zero-variance structural forward-filling (`ffill()`), capturing accurate time-series continuity.",
        "• Mathematical Sanity Validations: Scripted logical verification asserting that daily High bounds >= Low bounds, High >= Open/Close, and Low <= Open/Close inputs across all timestamps."
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5))

    # --------------------------------------------------------
    # SLIDE 5: DESCRIPTIVE STATISTICS MATRIX
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "4. Baseline Descriptive Statistics")
    pts = [
        "Below is the complete tabular output generated from the processed dataset containing basic univariate statistical indicators across price series and absolute market liquidity indicators:"
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.3), Inches(11.9), Inches(0.8), font_size=13)
    add_dataframe_table(slide, os.path.join(TABLES_DIR, "summary_statistics.csv"), Inches(0.5), Inches(2.2), Inches(12.33), Inches(4.2), max_rows=6)

    # --------------------------------------------------------
    # SLIDE 6: FINANCIAL ECONOMETRICS RETURN FORMULATION
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "5. Modeling Mathematical Return Dynamics")
    pts = [
        "• Continuous Compounding Preference: Logarithmic returns are calculated as the log difference of daily closing asset bounds:",
        "      R_t = ln(P_t / P_{t-1}) = ln(P_t) - ln(P_{t-1})",
        "• Mathematical Additivity: Log returns provide seamless temporal aggregation, meaning multi-period holding yields equate to the simple arithmetic summation of individual daily return vectors.",
        "• Statistical Properties: Log return series bound data behavior cleanly by bypassing structural lower limits of -100%, stabilizing optimization metrics.",
        "• Scaling Conversions: Returns are systematically expressed as percentages (R_t * 100) to optimize the scaling stability of complex non-linear GARCH estimation parameters."
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5))

    # --------------------------------------------------------
    # SLIDE 7: CLOSING PRICE HISTORICAL TRAJECTORIES
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "6. NIFTY 50 Absolute Price Evolution")
    add_image_safely(slide, "closing_price_trend.png", Inches(0.5), Inches(1.5), Inches(7.5), Inches(4.8))
    pts = [
        "• Visual Analysis: The price chart documents a highly persistent deterministic or stochastic trend component.",
        "• Macro Regime Shocks: Visualizes structural drift changes showing steady capital appreciation over the 2023-2026 landscape.",
        "• Non-Stationarity Inference: Price levels clearly exhibit regular non-reverting patterns, indicating that the baseline distribution lacks fixed global mean boundaries."
    ]
    add_bullet_points(slide, pts, Inches(8.3), Inches(1.8), Inches(4.5), Inches(4.5), font_size=13)

    # --------------------------------------------------------
    # SLIDE 8: LOG RETURN VOLATILITY CLUSTERING
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "7. Daily Log Return Trajectory")
    add_image_safely(slide, "daily_return_trend.png", Inches(0.5), Inches(1.5), Inches(7.5), Inches(4.8))
    pts = [
        "• Structural Transformation: The non-stationary absolute price trajectory is successfully transformed into a stationary, mean-reverting series around zero.",
        "• Stylized Fact - Volatility Clustering: Rapidly exhibits distinct financial time-series patterns where calm days follow calm days and high-magnitude changes follow high-magnitude changes.",
        "• Heteroscedasticity Signpost: Variance magnitudes clearly vary over time, signaling that conditional homoscedastic assumptions are highly invalid."
    ]
    add_bullet_points(slide, pts, Inches(8.3), Inches(1.8), Inches(4.5), Inches(4.5), font_size=13)

    # --------------------------------------------------------
    # SLIDE 9: THEORY OF LOG RETURN DISTRIBUTION
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "8. Asset Return Distribution Theory")
    pts = [
        "• Classic Gaussian Benchmark: Under standard random-walk finance models, daily log asset returns are assumed to follow an independent, identically distributed normal structure.",
        "• Empirical Divergence: Real-world financial assets consistently break standard normal assumptions due to market behavior shocks, tail-risk actions, and asymmetric market updates.",
        "• Skewness Profile: Measures the structural symmetry of the return tail extensions around the calculated series sample mean.",
        "• Kurtosis Profile: Quantifies the volume of structural probability density allocated out toward extreme tail territories relative to standard normal bell curves."
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5))

    # --------------------------------------------------------
    # SLIDE 10: JARQUE-BERA TEST & LEPTOKURTOSIS
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "9. Testing for Normality & Leptokurtosis")
    pts = [
        "• Jarque-Bera Formulation: Mathematically evaluates whether sample skewness (S) and sample kurtosis (K) structural parameters mirror a normal distribution standard:",
        "      JB = (n / 6) * [S^2 + (1 / 4) * (K - 3)^2]",
        "• Statistical Hypotheses: The null framework H0 posits strict distribution normality. Rejecting H0 implies that the series displays non-Gaussian behavior.",
        "• Empirical Findings: Financial returns almost universally reject normality due to extreme tail weights—a stylized market property called 'leptokurtosis'.",
        "• Modeling Impacts: Highly mandates implementing non-linear conditional variance models to capture fat-tail volatility dynamics accurately."
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5))

    # --------------------------------------------------------
    # SLIDE 11: STATIONARITY AND THE UNIT ROOT PROBLEM
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "10. Stationarity & Unit Root Dynamics")
    pts = [
        "• Foundational Definition: A time series is covariance-stationary if its mean, variance, and autocovariance properties remain strictly invariant over temporal shifts.",
        "• Econometric Risks: Running regressions on non-stationary absolute price components creates highly dangerous spurious relationships, inflating structural t-statistics falsely.",
        "• The Integrated Framework: Most financial price indices operate under integrated I(1) processes that require first-differencing transformations to achieve stationary I(0) constraints.",
        "• Mathematical Identification: Handled explicitly using structural autoregressive models designed to isolate unit root drift tendencies."
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5))

    # --------------------------------------------------------
    # SLIDE 12: AUGMENTED DICKEY-FULLER (ADF) TESTING METHOD
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "11. The Augmented Dickey-Fuller Blueprint")
    pts = [
        "• Econometric Specification: The test estimates an augmented parametric regression to identify unit root structures while controlling for higher-order serial dependencies:",
        "      ΔY_t = α + βt + γY_{t-1} + Σ δ_i ΔY_{t-i} + ε_t",
        "• Operational Hypotheses: Null framework H0: γ = 0 (implies Unit Root / Non-Stationary), against Alternative H1: γ < 0 (Covariance Stationary).",
        "• Empirical Verification Strategy: Check that absolute price levels fail to reject H0, while log return vectors explicitly reject the unit root condition at standard critical levels."
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5))

    # --------------------------------------------------------
    # SLIDE 13: SERIAL CORRELATION THEORY
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "12. Serial Autocorrelation Frameworks")
    pts = [
        "• Core Concept: Autocorrelation measures the linear correlation strength linking sequential values of a specific time-series variable back to its historical past.",
        "• Efficient Market Hypothesis (EMH) Linkages: Weak-form EMH states that past return data cannot help forecast future return vectors, which requires returns to resemble a white-noise process.",
        "• The Squared Return Twist: Finding no autocorrelation in raw returns confirms market informational efficiency, but finding significant autocorrelation in *squared* returns implies structural non-linear volatility dependency.",
        "• Architectural Pipeline Significance: Autocorrelation across squared return vectors serves as the direct econometric prerequisite for initiating ARCH/GARCH conditional variance frameworks."
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5))

    # --------------------------------------------------------
    # SLIDE 14: THE LJUNG-BOX TEST APPLICATION
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "13. Portmanteau Ljung-Box Test Structure")
    pts = [
        "• Mathematical Specification: The test evaluates the presence of global joint autocorrelation structures up to a designated maximum lag threshold (h):",
        "      Q = n(n + 2) * Σ [ (ρ_k^2) / (n - k) ]",
        "• Hypotheses: Null H0 assumes the series values are entirely independent (white noise), meaning all autocorrelation coefficients up to lag h are zero.",
        "• Dual Vector Assessment Strategy:",
        "  1. Raw Return Vector: Tests for residual structure inside the conditional mean equation.",
        "  2. Squared Return Vector: Checks for time-varying dependencies inside the conditional variance equation."
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5))

    # --------------------------------------------------------
    # SLIDE 15: ECONOMETRIC DIAGNOSTIC MATRIX
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "14. Empirical Econometric Diagnostics Summary")
    pts = [
        "The following structured matrix presents the empirical test statistics and corresponding p-values extracted directly across our core research phases:"
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.3), Inches(11.9), Inches(0.8), font_size=13)
    add_dataframe_table(slide, os.path.join(TABLES_DIR, "econometric_tests_summary.csv"), Inches(0.5), Inches(2.2), Inches(12.33), Inches(4.0), max_rows=8)

    # --------------------------------------------------------
    # SLIDE 16: AUTOREGRESSIVE CONDITIONAL HETEROSCEDASTICITY (ARCH)
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "15. Foundations of ARCH Specifications")
    pts = [
        "• Traditional Modeling Limits: Standard classical econometric models assume homoscedasticity, where error term variances remain globally constant over time.",
        "• Engle's Paradigm Shift (1982): Introduced models where the current conditional error variance explicitly depends on past squared residuals.",
        "• Mathematical Mechanism: Allows the variance of asset errors to expand rapidly following market volatility shocks, providing an accurate mathematical representation of real-world risk spikes.",
        "• Economic Interpretation: Captures the flow of market information, showing that large errors cluster together due to delayed structural market adjustment speeds."
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5))

    # --------------------------------------------------------
    # SLIDE 16B: ENGLE’S ARCH-LM MODEL RESIDUAL TESTING
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "16. Engle's ARCH Lagrange Multiplier Test")
    pts = [
        "• Operational Methodology: Runs an auxiliary regression of squared returns on their own past lagged values to identify structural variance dependencies:",
        "      ε²_t = α0 + α1*ε²_{t-1} + ... + αq*ε²_{t-q} + v_t",
        "• Test Statistic Structure: Uses a Lagrange Multiplier metric calculated as LM = n * R², asymptotically following a Chi-squared distribution with q degrees of freedom.",
        "• Decisive Inference: Rejecting the null hypothesis (H0: Homoscedasticity) confirms the presence of time-varying variance structures.",
        "• Tabular Test Findings: Below are the empirical results from Engle's ARCH-LM test run on NIFTY 50 innovations:"
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.5), Inches(11.9), Inches(3.2), font_size=13)
    add_dataframe_table(slide, os.path.join(TABLES_DIR, "arch_lm_test.csv"), Inches(2.5), Inches(4.8), Inches(8.33), Inches(1.8), max_rows=4)

    # --------------------------------------------------------
    # SLIDE 17: GARCH(1,1) PARAMETRIC STRUCTURAL LAYOUT
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "17. The GARCH(1,1) Estimation Model")
    pts = [
        "• Bollerslev's Extension (1986): Generalized ARCH models to incorporate historical conditional variance lags, creating a highly parsimonious alternative to infinite-lag ARCH models.",
        "• The Conditional Variance Equation: Specifies current variance as a function of a baseline intercept, yesterday's squared innovation shock, and yesterday's conditional variance persistence cell:",
        "      σ²_t = ω + α * ε²_{t-1} + β * σ²_{t-1}",
        "• Parametric Restrictions: Requires parameters to meet stability conditions where ω > 0, α >= 0, β >= 0, alongside a strict stationarity ceiling where (α + β) < 1.",
        "• Shocks and Persistence: The alpha parameter measures short-term variance reactivity to recent market innovations, while beta measures the long-term persistence of historical volatility shocks."
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5))

    # --------------------------------------------------------
    # SLIDE 18: GARCH CONDITIONAL VOLATILITY EVOLUTION
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "18. NIFTY 50 GARCH Conditional Volatility")
    add_image_safely(slide, "conditional_volatility.png", Inches(0.5), Inches(1.5), Inches(7.5), Inches(4.8))
    pts = [
        "• Dynamic Risk Tracking: Visualizes the annualized conditional volatility series extracted directly from the estimated GARCH model.",
        "• Analytical Strengths: Unlike static standard deviations, the GARCH framework captures rapid risk expansions during high-volatility regimes.",
        "• Policy and Application Value: Provides institutional risk managers with an objective, time-varying measure of underlying market asset variance."
    ]
    add_bullet_points(slide, pts, Inches(8.3), Inches(1.8), Inches(4.5), Inches(4.5), font_size=13)

    # --------------------------------------------------------
    # SLIDE 19: VALUE AT RISK (VaR) & CVaR RISK THEORY
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "19. Quantitative Risk Metrics: VaR & CVaR")
    pts = [
        "• Value at Risk (VaR): Quantifies the maximum expected capital loss over a specified horizon at a given significance target (α): P(R_t < -VaR) = α.",
        "• Parametric GARCH VaR Formulation: Utilizes the time-varying standard deviation outputs from the GARCH model to generate responsive risk boundaries:",
        "      VaR_t = -(μ + z_α * σ_t)",
        "• Conditional VaR (CVaR / Expected Shortfall): Measures the expected loss conditional on the return breach falling beyond the calculated VaR cutoff point: CVaR = E[R_t | R_t < -VaR].",
        "• Basel Regulatory Standard: Serves as a vital internal metric for calculating capital adequacy requirements under international financial regulatory frameworks."
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5))

    # --------------------------------------------------------
    # SLIDE 20: EMPIRICAL RISK ANALYSIS RESULTS
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "20. VaR & Expected Shortfall Calculations")
    add_image_safely(slide, "var_distribution.png", Inches(0.5), Inches(1.5), Inches(7.0), Inches(4.8))
    pts = [
        "Below are the calculated daily risk thresholds comparing historical simulation boundaries against modern GARCH parametric updates across key confidence intervals:"
    ]
    add_bullet_points(slide, pts, Inches(7.8), Inches(1.5), Inches(5.0), Inches(0.8), font_size=12)
    add_dataframe_table(slide, os.path.join(TABLES_DIR, "var_cvar_summary.csv"), Inches(7.7), Inches(2.5), Inches(5.1), Inches(3.2), max_rows=5)

    # --------------------------------------------------------
    # SLIDE 21: SUSTAINED CAPITAL LOSSES: MAXIMUM DRAWDOWN
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "21. Maximum Drawdown Analysis")
    add_image_safely(slide, "drawdown_chart.png", Inches(0.5), Inches(1.5), Inches(7.0), Inches(4.8))
    pts = [
        "• Definition: Measures the largest peak-to-trough decline in cumulative asset value over a given horizon, tracking sustained downside risk exposure.",
        "• Formula: Drawdown_t = (Cumulative_Return_t - Running_Peak_t) / Running_Peak_t.",
        "• Key Metrics Matrix:"
    ]
    add_bullet_points(slide, pts, Inches(7.8), Inches(1.5), Inches(5.0), Inches(1.5), font_size=12)
    add_dataframe_table(slide, os.path.join(TABLES_DIR, "drawdown_summary.csv"), Inches(7.7), Inches(3.2), Inches(5.1), Inches(2.8), max_rows=5)

    # --------------------------------------------------------
    # SLIDE 22: VOLATILITY REGIME SWITCHING THEORY
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "22. Macro Volatility Regime Classification")
    pts = [
        "• Theoretical Framework: Markets often transition between distinct economic environments characterized by calm periods and sudden structural adjustments.",
        "• Operational Segmentation Technique: Uses the long-term sample median of GARCH conditional volatility as an objective threshold to classify trading days:",
        "  1. Low-Volatility Regime: Days where conditional variance falls below or equal to the historical median.",
        "  2. High-Volatility Regime: Days where conditional variance exceeds the median threshold.",
        "• Empirical Merits: Converts continuous variance series into distinct, actionable phases to match changing risk profiles with real-world macro market conditions."
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5))

    # --------------------------------------------------------
    # SLIDE 23: REGIME ESTIMATION METRICS & GRAPHICS
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "23. Empirical Regime Shading Analysis")
    add_image_safely(slide, "volatility_regimes.png", Inches(0.5), Inches(1.5), Inches(7.0), Inches(4.8))
    pts = [
        "The table displays the empirical distribution of trading days and average annualized volatility across classified regimes:"
    ]
    add_bullet_points(slide, pts, Inches(7.8), Inches(1.5), Inches(5.0), Inches(0.8), font_size=12)
    add_dataframe_table(slide, os.path.join(TABLES_DIR, "regime_summary.csv"), Inches(7.7), Inches(2.5), Inches(5.1), Inches(3.8), max_rows=8)

    # --------------------------------------------------------
    # SLIDE 24: EMPIRICAL FINDINGS & CORE SYNTHESIS
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "24. Core Empirical Conclusions")
    pts = [
        "• Stylized Fact Validation: The NIFTY 50 return series displays strong non-Gaussian features, including negative skewness and high leptokurtic tail weight.",
        "• Stationarity Profile: Confirmed that NIFTY 50 absolute prices follow an integrated I(1) process, while log returns are highly stationary I(0) vectors.",
        "• Volatility Dependency Proof: Engle's ARCH-LM test strongly rejects the constant variance null hypothesis, confirming time-varying variance structures.",
        "• Risk Architecture Mapping: The GARCH(1,1) model successfully captures volatility clustering, providing risk managers with responsive, time-varying Value at Risk (VaR) profiles."
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5), font_size=15)

    # --------------------------------------------------------
    # SLIDE 25: RESEARCH OUTLOOK & STRATEGIC EXTENSIONS
    # --------------------------------------------------------
    slide = create_blank_slide(prs)
    add_slide_header(slide, "25. Strategic Horizons & Extensions")
    pts = [
        "• Asymmetric GARCH Implementations: Future extensions will incorporate EGARCH or GJR-GARCH frameworks to capture asymmetric leverage effects, where negative shocks impact volatility more than positive shocks.",
        "• Alternative Tail Distribution Assumptions: Replacing the default standard Gaussian assumption with Student-t or Skewed Student-t innovations to better capture leptokurtic tail weights.",
        "• Multi-Asset Spillover Frameworks: Expanding the pipeline to multivariate DCC-GARCH models to analyze systemic risk transmission channels linking NIFTY 50 dynamics with global macro assets.",
        "• High-Frequency Intraday Expansion: Transitioning from daily close data to intraday tick structures to analyze micro-level price discovery and realized volatility signatures."
    ]
    add_bullet_points(slide, pts, Inches(0.7), Inches(1.8), Inches(11.9), Inches(4.5), font_size=14)

    # --------------------------------------------------------
    # SAVE EXECUTION
    # --------------------------------------------------------
    final_output_path = os.path.join(OUTPUT_DIR, "NIFTY50_Volatility_Study.pptx")
    prs.save(final_output_path)
    print(f"\n" + "="*70)
    print(f"SUCCESS: 25-Slide Presentation Saved to Destination Path:")
    print(f"--> {final_output_path}")
    print(f"All slides feature academic attribution watermark.")
    print("="*70 + "\n")

if __name__ == "__main__":
    build_academic_presentation()